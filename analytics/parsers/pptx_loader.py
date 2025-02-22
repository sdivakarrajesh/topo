from analytics.parsers.base import FileParser
from pptx import Presentation
import re
from analytics.models import AnnualSummary, RevenueBreakdown, ActivityType, Company, QuarterlyPerformanceReport

class PPTXParser(FileParser):
    def __init__(self, file_path):
        self.file_path = file_path
        self.annual_summary = None
        self.revenue_breakdown = []
        self.quarterly_reports = []
        self.year = 2023 # Assuming it's a 2023 summary

    def parse(self):

        print(f"Parsing {self.file_path}")
        prs = Presentation(self.file_path)

        # slide number to parser function mapping
        parser_fns = {
            0: self.extract_key_highlights,     # Slide 1: Total Revenue, Memberships, Location
            1: self.extract_quarterly_metrics,  # Slide 2: Quarterly Revenue & Memberships
            2: self.extract_revenue_breakdown,  # Slide 3: Revenue Breakdown by Activity
        }

        for slide_num, slide in enumerate(prs.slides):
            print(f"Processing Slide {slide_num + 1}")
            if slide_num in parser_fns:
                parser_fns[slide_num](slide)

        print("Finished parsing PPTX.")

    def extract_key_highlights(self, slide):
        """
        Extracts total revenue, memberships sold, and top location from the first slide.
        """
        text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])

        revenue_match = re.search(r"Total Revenue: \$([\d,]+)", text)
        memberships_match = re.search(r"Total Memberships Sold: ([\d,]+)", text)
        location_match = re.search(r"Top Location: (.+)", text)

        revenue = float(revenue_match.group(1).replace(",", "")) if revenue_match else 0.0
        memberships_sold = int(memberships_match.group(1).replace(",", "")) if memberships_match else 0
        top_location = location_match.group(1) if location_match else "Unknown"

        company, _ = Company.objects.get_or_create(name="FitPro", defaults={"industry": "Sports and Leisure"})

        self.annual_summary = AnnualSummary(
            company=company,
            year=self.year,
            total_revenue=revenue,
            total_memberships_sold=memberships_sold,
            top_location=top_location
        )
        print(f"Extracted Key Highlights: Revenue={revenue}, Memberships={memberships_sold}, Location={top_location}")

    def extract_quarterly_metrics(self, slide):
        """
        Extracts quarterly revenue and memberships sold from a table in Slide 2.
        """
        for shape in slide.shapes:
            if hasattr(shape, "table"):  # Ensure the shape contains a table
                table = shape.table

                # Loop through table rows (skip header by using an index)
                for row_idx in range(1, len(table.rows)):
                    row = table.rows[row_idx]
                    try:
                        year = self.year
                        quarter = row.cells[0].text.strip()
                        revenue = float(row.cells[1].text.strip().replace(",", "").replace("$", ""))
                        memberships_sold = int(row.cells[2].text.strip())
                        avg_duration = int(row.cells[3].text.strip())

                        self.quarterly_reports.append(QuarterlyPerformanceReport(
                            report_type="FitPro: Annual Summary 2023",
                            year=year,
                            quarter=quarter,
                            revenue=revenue,
                            memberships_sold=memberships_sold,
                            avg_duration=avg_duration
                        ))
                        print(f"Extracted Quarterly Metrics: {year} {quarter} | Revenue: {revenue} | Memberships: {memberships_sold} | Avg Duration: {avg_duration}")
                    except ValueError as e:
                        print(f"Skipping invalid row {row_idx}: {e}")


    def extract_revenue_breakdown(self, slide):
        """
        Extracts revenue distribution per activity from the third slide.
        """
        text_sections = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        while text_sections:
            section = text_sections.pop(0)
            if "Revenue Distribution:" in section:
                text_lines = section.split("\n")[1:]
                break

        if text_lines:
            for activity_data in text_lines:
                parts = activity_data.split(":")
                if len(parts) == 2:
                    activity_name = parts[0].strip()
                    percentage = float(parts[1].replace("%", "").strip())

                    activity_type, _ = ActivityType.objects.get_or_create(name=activity_name)

                    self.revenue_breakdown.append(RevenueBreakdown(
                        annual_summary=self.annual_summary,
                        activity_type=activity_type,
                        percentage=percentage
                    ))

            print(f"Extracted Revenue Breakdown: {len(self.revenue_breakdown)} entries.")
        else:
            print("No revenue breakdown found in the slide.")

    def load_db(self):
        """
        Bulk insert the parsed records into the database.
        """
        print(f"Loading {self.file_path} data into the database")

        if self.annual_summary:
            self.annual_summary.save()

            # Attach the saved summary to revenue breakdowns and bulk create
            for breakdown in self.revenue_breakdown:
                breakdown.annual_summary = self.annual_summary
            RevenueBreakdown.objects.bulk_create(self.revenue_breakdown)

        # Bulk insert quarterly performance reports
        if self.quarterly_reports:
            QuarterlyPerformanceReport.objects.bulk_create(self.quarterly_reports)

        print("Insertion done!")

    def run(self):
        self.parse()
        self.load_db()
